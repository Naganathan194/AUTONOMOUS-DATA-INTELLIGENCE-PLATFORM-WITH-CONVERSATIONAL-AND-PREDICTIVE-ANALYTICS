import re
import logging
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import datetime

from utils import get_gemini_response

def clean_ai_text(text: str) -> str:
    """
    Comprehensive cleaning of AI-generated text:
    - Removes control characters
    - Removes carriage returns
    - Removes asterisks (*) and backticks (`)
    - Removes markdown-style headers (# ## ###)
    - Collapses multiple blank lines
    - Cleans bullet point formatting
    """
    if not text:
        return ""
    
    # 1. Remove ASCII control characters
    cleaned = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F-\x9F]', '', text)
    
    # 2. Remove carriage returns
    cleaned = cleaned.replace('\r', '')
    
    # 3. Remove all asterisks and backticks
    cleaned = re.sub(r'[*`]', '', cleaned)
    
    # 4. Remove markdown headers (# ## ###)
    cleaned = re.sub(r'^#{1,6}\s+', '', cleaned, flags=re.MULTILINE)
    
    # 5. Remove markdown bold/italic markers (**text** or __text__)
    cleaned = re.sub(r'\*\*(.+?)\*\*', r'\1', cleaned)
    cleaned = re.sub(r'__(.+?)__', r'\1', cleaned)
    
    # 6. Clean up bullet points - normalize to "• "
    cleaned = re.sub(r'^[\s]*[-\*\+]\s+', '• ', cleaned, flags=re.MULTILINE)
    
    # 7. Collapse multiple blank lines into double newline
    cleaned = re.sub(r'\n\s*\n\s*\n+', '\n\n', cleaned)
    
    # 8. Remove extra spaces around bullet points
    cleaned = re.sub(r'\n\s+•', '\n•', cleaned)
    
    # 9. Trim leading/trailing whitespace
    cleaned = cleaned.strip()
    
    return cleaned

def to_bullets(text: str, max_bullets: int = 10) -> str:
    """
    Convert text into clean bullet list format:
    - Preserves existing bullet points
    - Splits paragraphs into bullets if needed
    - Limits to max_bullets
    - Ensures consistent formatting
    """
    if not text:
        return ""
    
    # Clean the text first
    text = clean_ai_text(text)
    
    # Split by newlines to get potential bullet items
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    
    bullets = []
    for line in lines:
        # Skip if already at max
        if len(bullets) >= max_bullets:
            break
        
        # Remove existing bullet markers
        line = re.sub(r'^[•\-\*\+]\s*', '', line)
        
        # Split long paragraphs by sentences if needed
        if len(line) > 200 and '.' in line:
            sentences = re.split(r'(?<=[.!?])\s+', line)
            for sent in sentences:
                sent = sent.strip()
                if sent and len(bullets) < max_bullets:
                    bullets.append(sent)
        elif line:
            bullets.append(line)
    
    # Return formatted bullets
    if bullets:
        return '• ' + '\n• '.join(bullets)
    return ""

def derive_basic_insights(eda_metadata, df):
    """Fallback insights if AI fails"""
    bullets = []
    try:
        num_rows, num_cols = df.shape
        bullets.append(f"Dataset contains {num_rows:,} rows and {num_cols} columns")
        
        cols = eda_metadata.get("columns", {})
        if cols:
            numeric_cols = [c for c, d in cols.items() if "numeric_stats" in d]
            categorical_cols = [c for c, d in cols.items() if d.get("dtype", "").lower() == "object"]
            
            if numeric_cols:
                bullets.append(f"Found {len(numeric_cols)} numeric columns for quantitative analysis")
            if categorical_cols:
                bullets.append(f"Found {len(categorical_cols)} categorical columns for classification")
            
            # Missing data summary
            missing = eda_metadata.get("missing_data_overall", {})
            if missing:
                high_missing = {k: v for k, v in missing.items() if v > 10}
                if high_missing:
                    bullets.append(f"Columns with significant missing data: {', '.join(high_missing.keys())}")
        
    except Exception:
        bullets.append("Dataset successfully loaded and analyzed")
    
    return '• ' + '\n• '.join(bullets)

def generate_eda_report_ppt(eda_metadata, df, dataset_name="EDA_Report.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_mpl_fig_to_slide(prs, fig, title):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        out = io.BytesIO()
        fig.savefig(out, format='png', bbox_inches='tight', dpi=170)
        plt.close(fig)
        out.seek(0)
        slide.shapes.add_picture(out, Inches(1), Inches(1.5), Inches(8), Inches(4))

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EDA Report"
    slide.placeholders[1].text = f"Dataset: {dataset_name}\nRows: {df.shape[0]}, Columns: {df.shape[1]}"

    # --- AI Insights Slides (Split into two) ---
    try:
        import json as json_module
        prompt = f"""Analyze the following dataset EDA and provide a key insights summary in 10-15 bullets.\n\nDataset: {dataset_name}\nRows: {len(df):,}, Columns: {len(df.columns)}\n\nEDA Summary:\n{json_module.dumps(eda_metadata, indent=2, default=str)[:2000]}\n\nSummarize actionable and interesting patterns, data quality, notable columns, trends, correlations, and anything unusual. Format as plain text bullets starting with '- '."""
        ai_text = get_gemini_response(prompt, "flash")
        bullets = to_bullets(ai_text, max_bullets=15)
        if not bullets:
            bullets = derive_basic_insights(eda_metadata, df)
    except Exception as e:
        logger = logging.getLogger(__name__)
        logger.warning(f"AI insights generation failed: {str(e)}")
        bullets = derive_basic_insights(eda_metadata, df)

    # Split bullets into two groups
    bullet_lines = [b.strip() for b in bullets.split('\n') if b.strip()]
    mid_point = len(bullet_lines) // 2 + (1 if len(bullet_lines) % 2 == 1 else 0)
    first_half = bullet_lines[:mid_point]
    second_half = bullet_lines[mid_point:]
    
    # Helper function to create insight slide
    def create_insight_slide(title_text, bullet_list):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)
        
        for idx, bullet_text in enumerate(bullet_list):
            if idx == 0:
                p = tf.paragraphs[0]
                p.text = bullet_text
            else:
                p = tf.add_paragraph()
                p.text = bullet_text
            
            p.font.size = Pt(20)
            p.level = 0
            p.space_after = Pt(8)
            if idx == 0:
                p.font.bold = True
    
    # Create first insights slide
    if first_half:
        create_insight_slide("Key Insights (AI-Powered) - Part 1", first_half)
    
    # Create second insights slide
    if second_half:
        create_insight_slide("Key Insights (AI-Powered) - Part 2", second_half)

    # --- Numeric (Histogram) Slides ---
    for col in df.select_dtypes(include='number').columns[:8]:
        fig, ax = plt.subplots(figsize=(8, 4))
        sns.histplot(df[col].dropna(), kde=True, color="#6C63FF", ax=ax)
        ax.set_title(f"{col} Distribution")
        ax.set_xlabel(col)
        ax.set_ylabel("Count")
        add_mpl_fig_to_slide(prs, fig, f"Distribution: {col}")

    # --- Categorical (Bar) Slides ---
    for col in df.select_dtypes(include=['object', 'category']).columns[:8]:
        vc = df[col].value_counts().head(12)
        fig, ax = plt.subplots(figsize=(8, 4))
        sns.barplot(x=vc.index, y=vc.values, ax=ax, palette="mako")
        ax.set_title(f"{col} (Top 12)")
        plt.xticks(rotation=30, ha='right')
        ax.set_ylabel("Count")
        ax.set_xlabel(col)
        add_mpl_fig_to_slide(prs, fig, f"Categorical: {col}")

    # --- Correlation Heatmap ---
    nums = df.select_dtypes(include="number")
    if nums.shape[1] >= 2:
        corr = nums.corr().fillna(0)
        fig, ax = plt.subplots(figsize=(7, 6))
        sns.heatmap(corr, annot=True, cmap="RdBu_r", center=0, ax=ax)
        ax.set_title("Correlation Heatmap")
        add_mpl_fig_to_slide(prs, fig, "Correlation Heatmap")

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer